from selenium import webdriver
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.common.action_chains import ActionChains
from PIL import Image
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from time import sleep
from glob import glob
import os
import argparse

class HandleYT:
    def __init__(self, url, temp_dir='tmp'):
        self.driver = driver = webdriver.Chrome("./chromedriver")
        self.filecounter = 0
        self.temp_dir = 'tmp'
        if not os.path.exists(self.temp_dir):
            os.makedirs(self.temp_dir)
        driver.get(url)

    def _pause(self):
        self.driver.execute_script("document.getElementsByTagName('video')[0].pause()")

    def _beginning(self):
        actions = ActionChains(self.driver)
        actions.send_keys(Keys.HOME)
        actions.perform()

    def _on_ad(self):
        return self.driver.execute_script("return document.getElementById('movie_player').getPlayerState()") == -1

    def _over_duration(self, duration):
        if duration:
            return self.driver.execute_script("return document.getElementsByTagName('video')[0].getCurrentTime()") >= duration
        else:
            return False

    def _next_frame(self):
        actions = ActionChains(self.driver)
        actions.send_keys('.')
        #actions.send_keys(Keys.ARROW_RIGHT)
        actions.perform()

    def _at_end(self):
        return self.driver.execute_script("return document.getElementById('movie_player').getPlayerState()") == 0

    def _loading(self):
        return not "display: none" in self.driver.find_elements_by_class_name("ytp-spinner")[0].get_attribute("style")

    def _take_screenshot(self):
        video = self.driver.find_element_by_id("player-container").screenshot_as_png
        im = Image.open(BytesIO(video)) # uses PIL library to open image in memory
        im.save(os.path.join(self.temp_dir, 'screenshot%d.png' % (self.filecounter,))) # saves new cropped image
        self.filecounter += 1

    def get_screenshots(self, duration=None):
        sleep(3)
        while self._on_ad():
            sleep(0.5)
        sleep(3)
        self._pause()
        self._beginning()
        print("Setup")
        while not self._at_end():
            while self._loading():
                sleep(0.02)
            if self._over_duration(duration):
                break
            self._take_screenshot()
            self._next_frame()
        self.driver.close()

class HandlePPT:
    def __init__(self):
        self.pres = Presentation()

    def add_picture(self, filename):
        blank_slide_layout = self.pres.slide_layouts[6]
        slide = self.pres.slides.add_slide(blank_slide_layout)
        left = top = Inches(0)
        pic = slide.shapes.add_picture(filename, left, top, self.pres.slide_width, self.pres.slide_height)

    def save(self, output):
        self.pres.save(output)

def main():
    ytdl = HandleYT(args.url, args.temp_dir)
    ytdl.get_screenshots(duration=args.duration)

    ppt = HandlePPT()
    for screenshot in glob('tmp/screenshot*.png'):
        ppt.add_picture(screenshot)

    ppt.save(args.output)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Download some youtube videos the long way.')
    parser.add_argument('url', help='URL to download.')
    parser.add_argument('output', help="Output pptx file")
    parser.add_argument('-d', "--duration", type=float, default=None, help='Optional seconds to download. If not specified, download everything')
    parser.add_argument('-t', '--temp-dir', default='tmp', help="Temporary directory to hold the screenshots")
    args = parser.parse_args()

    main()